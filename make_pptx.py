# -*- coding: utf-8 -*-
"""Erstellt die PowerPoint-Praesentation 'Werkstoff Aluminium'."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ---- Farbpalette -----------------------------------------------------------
BLUE_D = RGBColor(0x0B, 0x3D, 0x91)
BLUE   = RGBColor(0x1F, 0x6F, 0xEB)
ALU    = RGBColor(0x9A, 0xA7, 0xB4)
ALU_D  = RGBColor(0x5C, 0x6B, 0x7A)
DARK   = RGBColor(0x22, 0x2A, 0x33)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x9E, 0x5B)
ORANGE = RGBColor(0xE0, 0x8A, 0x1E)
RED    = RGBColor(0xC0, 0x32, 0x32)
GREY_T = RGBColor(0x66, 0x70, 0x7A)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri", line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def bullets(s, x, y, w, h, items, size=18, color=DARK, gap=6, bullet_color=BLUE,
            bold_lead=False, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        # Bullet-Symbol
        rb = p.add_run(); rb.text = "▸  "
        rb.font.size = Pt(size); rb.font.color.rgb = bullet_color; rb.font.bold = True
        rb.font.name = "Calibri"
        if isinstance(it, tuple):  # (lead, rest)
            r1 = p.add_run(); r1.text = it[0]
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name="Calibri"
            r2 = p.add_run(); r2.text = it[1]
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name="Calibri"
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name="Calibri"
            r.font.bold = bold_lead
    return tb

def header(s, num, title, accent=BLUE):
    """Kopfzeile mit Aufgaben-Nummer-Badge."""
    rect(s, 0, 0, SW, Inches(1.15), BLUE_D)
    rect(s, 0, Inches(1.15), SW, Pt(5), accent)
    if num:
        badge = rect(s, Inches(0.45), Inches(0.27), Inches(0.62), Inches(0.62),
                     accent, shape=MSO_SHAPE.OVAL)
        tf = badge.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(num)
        r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name="Calibri"
        tx = Inches(1.25)
    else:
        tx = Inches(0.5)
    txt(s, tx, Inches(0.1), SW-tx-Inches(0.3), Inches(1.0), title,
        size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

def source(s, text):
    txt(s, Inches(0.45), SH-Inches(0.42), SW-Inches(0.9), Inches(0.35),
        "Quelle: " + text, size=10.5, color=GREY_T, italic=True)

def pageno(s, n):
    txt(s, SW-Inches(0.8), SH-Inches(0.42), Inches(0.5), Inches(0.35),
        str(n), size=11, color=GREY_T, align=PP_ALIGN.RIGHT)

def add_img(s, path, x, y, w=None, h=None, frame=True):
    iw, ih = Image.open(path).size
    ar = iw/ih
    if w and not h: h = Emu(int(int(w)/ar))
    if h and not w: w = Emu(int(int(h)*ar))
    if frame:
        rect(s, x-Pt(3), y-Pt(3), w+Pt(6), h+Pt(6), WHITE, line=ALU)
    s.shapes.add_picture(path, x, y, w, h)
    return w, h

A = "assets/"

# ===========================================================================
# Folie 1 – Titel
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, BLUE_D)
rect(s, 0, 0, Inches(0.35), SH, BLUE)
rect(s, 0, Inches(4.9), SW, Inches(2.6), RGBColor(0x09, 0x2F, 0x6E))
txt(s, Inches(0.9), Inches(0.8), Inches(8.5), Inches(0.5),
    "BSZ für Technik · Lernfeld 5 · Nutzfahrzeugbau", size=16, color=ALU, bold=True)
txt(s, Inches(0.85), Inches(1.45), Inches(11), Inches(1.6),
    "Der Werkstoff Aluminium", size=54, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(3.0), Inches(11), Inches(1.1),
    "Vorkommen · Herstellung · Legierungen · Anwendung\nam Beispiel des Mercedes-Benz 2440 L",
    size=22, color=LIGHT, line_spacing=1.15)
w,h = add_img(s, A+"titel.png", Inches(0.9), Inches(4.95), h=Inches(2.2), frame=False)
txt(s, Inches(8.7), Inches(5.2), Inches(4.2), Inches(1.8),
    "Handout / Präsentation\nLehrjahr 2 – Werkstoffkunde\n12.06.2026",
    size=15, color=ALU, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# Folie 2 – Gliederung / Agenda
# ===========================================================================
s = slide()
header(s, None, "Gliederung")
items = [
    ("1   Vorkommen & Abbaugebiete", "– Wo findet man Aluminium auf der Erde?"),
    ("2   Vor- und Nachteile", "– Eigenschaften des Werkstoffs"),
    ("3   Herstellung", "– Vom Bauxit zum Reinaluminium"),
    ("4   Legierungen", "– aushärtbar/nicht aushärtbar · Knet- & Gusslegierung"),
    ("5   Festigkeitssteigerung", "– drei wichtige Verfahren"),
    ("6   Anwendung im Berufsalltag", "– Beispiele im Nutzfahrzeugbau"),
]
y = Inches(1.55)
for i, (lead, rest) in enumerate(items):
    card = rect(s, Inches(0.7), y, Inches(11.9), Inches(0.78), LIGHT)
    rect(s, Inches(0.7), y, Inches(0.14), Inches(0.78), BLUE)
    txt(s, Inches(1.1), y, Inches(4.6), Inches(0.78), lead, size=19, bold=True,
        color=BLUE_D, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.5), y, Inches(7.0), Inches(0.78), rest, size=15, color=ALU_D,
        anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.92)
pageno(s, 2)

# ===========================================================================
# Folie 3 – Aufgabe 1: Vorkommen & Abbaugebiete
# ===========================================================================
s = slide()
header(s, 1, "Vorkommen & Abbaugebiete")
bullets(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(4.6), [
    ("Häufigstes Metall: ", "Aluminium macht ca. 8 % der Erdkruste aus – Platz 3 aller Elemente (nach O und Si)."),
    ("Nie gediegen: ", "kommt wegen seiner Reaktionsfreudigkeit nur gebunden in Mineralien vor."),
    ("Wichtigstes Erz – Bauxit: ", "enthält 50–60 % Aluminiumoxid (Al₂O₃) sowie Eisen- und Siliciumoxide."),
    ("Abbau im Tagebau: ", "Bauxit liegt oft oberflächennah in tropischen/subtropischen Gebieten."),
    ("Lagerstätten weltweit: ", "Australien, Guinea, China, Brasilien, Indien, Indonesien, Jamaika."),
], size=16.5, gap=11)
add_img(s, A+"vorkommen.png", Inches(7.7), Inches(1.5), w=Inches(4.85))
add_img(s, A+"abbau.png", Inches(7.7), Inches(4.35), w=Inches(4.85))
source(s, "Bundesanstalt für Geowissenschaften und Rohstoffe (BGR); Roloff/Matek Tabellenbuch Metall; Wikipedia 'Aluminium'/'Bauxit' (2024).")
pageno(s, 3)

# ===========================================================================
# Folie 4 – Aufgabe 2: Vorteile / Nachteile
# ===========================================================================
s = slide()
header(s, 2, "Vor- und Nachteile des Werkstoffs", accent=GREEN)
# Vorteile
rect(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.6), GREEN)
txt(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(0.6), "✓  VORTEILE", size=20,
    color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.3), LIGHT)
bullets(s, Inches(0.8), Inches(2.25), Inches(5.7), Inches(4.1), [
    ("Geringe Dichte ", "(2,7 g/cm³) – nur ⅓ von Stahl → Leichtbau"),
    ("Hohe Korrosionsbeständigkeit ", "durch dichte Oxidschicht"),
    ("Gut umformbar ", "– walzen, pressen, tiefziehen, schweißen"),
    ("Sehr gute Leitfähigkeit ", "für Wärme und Strom"),
    ("Unmagnetisch ", "und lebensmittelecht"),
    ("Sehr gut recycelbar ", "– nur ~5 % der Primärenergie"),
], size=15.5, gap=9, bullet_color=GREEN)
# Nachteile
rect(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.6), RED)
txt(s, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.6), "✗  NACHTEILE", size=20,
    color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(6.9), Inches(2.1), Inches(6.0), Inches(4.3), LIGHT)
bullets(s, Inches(7.1), Inches(2.25), Inches(5.7), Inches(4.1), [
    ("Geringere Festigkeit ", "als Stahl (rein); E-Modul nur ⅓"),
    ("Hoher Energiebedarf ", "bei der Primärherstellung (Elektrolyse)"),
    ("Teurer ", "als Baustahl"),
    ("Schwerer schweißbar ", "– Oxidschicht, hohe Wärmeleitung"),
    ("Geringe Warmfestigkeit ", "– Festigkeitsverlust bei Hitze"),
    ("Niedrige Dauerfestigkeit ", "– keine echte Dauerfestigkeitsgrenze"),
], size=15.5, gap=9, bullet_color=RED)
source(s, "Roloff/Matek Maschinenelemente; Europa-Lehrmittel 'Fachkunde Metall' (58. Aufl.).")
pageno(s, 4)

# ===========================================================================
# Folie 5 – Aufgabe 3: Herstellung
# ===========================================================================
s = slide()
header(s, 3, "Herstellung – vom Bauxit zum Aluminium", accent=ORANGE)
add_img(s, A+"herstellung.png", Inches(0.55), Inches(1.4), w=Inches(8.4))
rect(s, Inches(9.2), Inches(1.5), Inches(3.7), Inches(4.7), LIGHT)
txt(s, Inches(9.4), Inches(1.6), Inches(3.4), Inches(0.5), "Kurz erklärt", size=18,
    bold=True, color=BLUE_D)
bullets(s, Inches(9.35), Inches(2.15), Inches(3.45), Inches(4.0), [
    ("Schritt 1: ", "Aus Bauxit wird im Bayer-Verfahren reine Tonerde (Al₂O₃) gewonnen."),
    ("Schritt 2: ", "Schmelzflusselektrolyse (Hall-Héroult): Al₂O₃ in Kryolith bei ~950 °C."),
    ("Strom spaltet ", "das Oxid – flüssiges Alu sammelt sich am Boden."),
    ("Energie: ", "~13–15 kWh je kg Rohaluminium."),
], size=13.5, gap=8, bullet_color=ORANGE)
source(s, "Europa-Lehrmittel 'Fachkunde Metall'; aluminiumdeutschland.de; Wikipedia 'Aluminium­herstellung' (2024).")
pageno(s, 5)

# ===========================================================================
# Folie 6 – Aufgabe 4: Legierungen Klassifikation
# ===========================================================================
s = slide()
header(s, 4, "Aluminiumlegierungen – Übersicht", accent=BLUE)
add_img(s, A+"legierungen.png", Inches(0.55), Inches(1.4), w=Inches(7.5))
rect(s, Inches(8.4), Inches(1.5), Inches(4.5), Inches(4.8), LIGHT)
txt(s, Inches(8.6), Inches(1.6), Inches(4.1), Inches(0.5),
    "aushärtbar vs. nicht aushärtbar", size=16, bold=True, color=BLUE_D)
bullets(s, Inches(8.55), Inches(2.2), Inches(4.2), Inches(4.0), [
    ("Aushärtbar: ", "Festigkeit lässt sich durch eine Wärmebehandlung (Ausscheidungshärtung) deutlich steigern – z. B. AlCuMg, AlMgSi."),
    ("Nicht aushärtbar: ", "keine Festigkeitssteigerung durch Wärmebehandlung möglich – nur durch Kaltverfestigung oder Legieren – z. B. AlMg, AlMn."),
], size=14, gap=12, bullet_color=BLUE)
source(s, "DIN EN 573 / EN 1706; Roloff/Matek Tabellenbuch; Europa-Lehrmittel 'Fachkunde Metall'.")
pageno(s, 6)

# ===========================================================================
# Folie 7 – Aufgabe 4b: Knet- vs Gusslegierung (Vergleich)
# ===========================================================================
s = slide()
header(s, 4, "Knetlegierung vs. Gusslegierung", accent=BLUE)
# Zwei Spalten
def vergleichskarte(x, titel, farbe, items):
    rect(s, x, Inches(1.5), Inches(5.95), Inches(0.65), farbe)
    txt(s, x, Inches(1.5), Inches(5.95), Inches(0.65), titel, size=20, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, Inches(2.15), Inches(5.95), Inches(4.4), LIGHT)
    bullets(s, x+Inches(0.2), Inches(2.3), Inches(5.6), Inches(4.2), items,
            size=15, gap=10, bullet_color=farbe)
vergleichskarte(Inches(0.6), "KNETLEGIERUNG", BLUE, [
    ("Verarbeitung: ", "durch Umformen – Walzen, Pressen, Schmieden, Ziehen."),
    ("Geringer Legierungsanteil ", "→ gut verformbar, zäh."),
    ("Halbzeuge: ", "Bleche, Profile, Rohre, Drähte."),
    ("Beispiele: ", "AlMg3, AlMgSi0,5, AlCuMg (Dural)."),
    ("Einsatz NFZ: ", "Aufbauten, Bordwände, Kofferplanken."),
])
vergleichskarte(Inches(6.78), "GUSSLEGIERUNG", ALU_D, [
    ("Verarbeitung: ", "durch Vergießen – Sandguss, Kokillen-, Druckguss."),
    ("Hoher Legierungsanteil ", "(v. a. Silicium) → gut gießbar, dünnflüssig."),
    ("Gute Formfüllung ", "auch komplexer Bauteile."),
    ("Beispiele: ", "AlSi12, AlSi10Mg, G-AlSi."),
    ("Einsatz NFZ: ", "Getriebe-/Motorgehäuse, Halter, Konsolen."),
])
source(s, "DIN EN 1706 (Guss) / EN 573 (Knet); Europa-Lehrmittel 'Fachkunde Metall'.")
pageno(s, 7)

# ===========================================================================
# Folie 8 – Aufgabe 5: Festigkeitssteigerung
# ===========================================================================
s = slide()
header(s, 5, "Festigkeitssteigerung – 3 Möglichkeiten", accent=GREEN)
add_img(s, A+"festigkeit.png", Inches(0.7), Inches(1.4), w=Inches(11.9))
# drei Erklaerkarten
cards = [
    (Inches(0.6), BLUE, "1  Mischkristallverfestigung",
     "Fremdatome (z. B. Mg, Mn) werden ins Gitter eingebaut. Sie verzerren das Kristallgitter und behindern das Abgleiten → höhere Festigkeit."),
    (Inches(4.62), ORANGE, "2  Kaltverfestigung",
     "Kaltumformen (Walzen, Ziehen) erhöht die Versetzungsdichte. Versetzungen stauen sich → Werkstoff wird fester, aber spröder."),
    (Inches(8.64), GREEN, "3  Ausscheidungshärtung",
     "Aushärten: Lösungsglühen → Abschrecken → Auslagern. Feine Ausscheidungen blockieren Versetzungen (nur bei aushärtbaren Legierungen)."),
]
for x, c, t, d in cards:
    rect(s, x, Inches(4.55), Inches(3.95), Inches(0.55), c)
    txt(s, x, Inches(4.55), Inches(3.95), Inches(0.55), t, size=15, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, Inches(5.1), Inches(3.95), Inches(1.7), LIGHT)
    txt(s, x+Inches(0.15), Inches(5.2), Inches(3.65), Inches(1.55), d, size=13,
        color=DARK, line_spacing=1.05)
source(s, "Bargel/Schulze 'Werkstoffkunde'; Europa-Lehrmittel 'Fachkunde Metall'; Roloff/Matek.")
pageno(s, 8)

# ===========================================================================
# Folie 9 – Aufgabe 6: Anwendung im Berufsalltag
# ===========================================================================
s = slide()
header(s, 6, "Anwendung im Berufsalltag (Nutzfahrzeug)", accent=ORANGE)
apps = [
    ("Fahrerhaus & Aufbau", "Bordwände, Kofferaufbauten und Planenrahmen aus Alu-Profilen senken das Leergewicht → höhere Nutzlast.", BLUE),
    ("Antrieb & Fahrwerk", "Motor- und Getriebegehäuse, Ölwannen, Querträger und Felgen als Aluminium-Gussteile.", ALU_D),
    ("Tanks & Behälter", "Kraftstoff- und Druckluftbehälter, Trittstufen und Unterfahrschutz aus korrosionsbeständigem Aluminium.", GREEN),
]
x = Inches(0.6)
for t, d, c in apps:
    rect(s, x, Inches(1.55), Inches(3.95), Inches(4.6), LIGHT)
    rect(s, x, Inches(1.55), Inches(3.95), Inches(0.95), c)
    txt(s, x+Inches(0.15), Inches(1.55), Inches(3.65), Inches(0.95), t, size=19,
        color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.25), Inches(2.7), Inches(3.5), Inches(3.2), d, size=16,
        color=DARK, line_spacing=1.2)
    x += Inches(4.18)
rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.62), BLUE_D)
txt(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.62),
    "Kernnutzen im NFZ-Bau:  weniger Gewicht  →  mehr Nutzlast, weniger Kraftstoff, weniger CO₂",
    size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
source(s, "Mercedes-Benz Trucks Technikinfo; aluminiumdeutschland.de; Europa-Lehrmittel 'Fachkunde Kfz-Technik'.")
pageno(s, 9)

# ===========================================================================
# Folie 10 – Zusammenfassung
# ===========================================================================
s = slide()
header(s, None, "Zusammenfassung")
items = [
    "Aluminium ist das häufigste Metall der Erdkruste, wird aber nur gebunden als Bauxit abgebaut.",
    "Größter Vorteil: geringe Dichte (2,7 g/cm³) und Korrosionsbeständigkeit → idealer Leichtbauwerkstoff.",
    "Herstellung in zwei Stufen: Bayer-Verfahren (Tonerde) + Schmelzflusselektrolyse – sehr energieintensiv, aber gut recycelbar.",
    "Knetlegierungen werden umgeformt, Gusslegierungen vergossen; nur aushärtbare Legierungen lassen sich wärmebehandeln.",
    "Festigkeit steigerbar durch Mischkristallverfestigung, Kaltverfestigung und Ausscheidungshärtung.",
    "Im Nutzfahrzeugbau spart Aluminium Gewicht und erhöht so die Nutzlast – am Mercedes-Benz 2440 L vielfach eingesetzt.",
]
y = Inches(1.55)
for i, it in enumerate(items):
    rect(s, Inches(0.7), y, Inches(0.55), Inches(0.55), GREEN, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(0.7), y, Inches(0.55), Inches(0.55), "✓", size=18, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.45), y-Inches(0.05), Inches(11.3), Inches(0.85), it, size=16.5,
        color=DARK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    y += Inches(0.9)
pageno(s, 10)

# ===========================================================================
# Folie 11 – Quellenangaben
# ===========================================================================
s = slide()
header(s, None, "Quellenangaben")
rect(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.4), LIGHT)
quellen = [
    ("Fachbücher", [
        "Europa-Lehrmittel: Fachkunde Metall, 58. Auflage, Verlag Europa-Lehrmittel, Haan-Gruiten.",
        "Roloff/Matek: Maschinenelemente – Tabellenbuch, Springer Vieweg.",
        "Bargel/Schulze: Werkstoffkunde, Springer Verlag.",
    ]),
    ("Normen", [
        "DIN EN 573 – Aluminium und Aluminiumlegierungen (Knetwerkstoffe).",
        "DIN EN 1706 – Aluminium und Aluminiumlegierungen (Gussstücke).",
    ]),
    ("Internet", [
        "aluminiumdeutschland.de – Aluminium Deutschland e. V. (abgerufen 06/2026).",
        "bgr.bund.de – Bundesanstalt für Geowissenschaften und Rohstoffe (Rohstoffdaten).",
        "de.wikipedia.org – Artikel 'Aluminium', 'Bauxit', 'Aluminiumherstellung' (abgerufen 06/2026).",
        "Mercedes-Benz Trucks – technische Produktinformationen.",
    ]),
    ("Aufgabenstellung / Bild", [
        "Aufgabenblatt 'Aluminium' (LF 5) sowie Fahrzeugabbildung Mercedes-Benz 2440 L.",
        "Grafiken/Diagramme: Eigene Darstellungen (schematisch).",
    ]),
]
y = Inches(1.7)
for kat, lst in quellen:
    txt(s, Inches(1.0), y, Inches(11.3), Inches(0.4), kat, size=16, bold=True, color=BLUE_D)
    y += Inches(0.42)
    for q in lst:
        txt(s, Inches(1.3), y, Inches(11.0), Inches(0.36), "•  " + q, size=12.5, color=DARK)
        y += Inches(0.33)
    y += Inches(0.08)
pageno(s, 11)

os.makedirs("output", exist_ok=True)
out = "Praesentation_Aluminium.pptx"
prs.save(out)
print("gespeichert:", out, "-", len(prs.slides._sldIdLst), "Folien")
